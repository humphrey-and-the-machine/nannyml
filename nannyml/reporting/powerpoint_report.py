#  Author:   Niels Nuyttens  <niels@nannyml.com>
#
#  License: Apache Software License 2.0
import re
from datetime import datetime
from pathlib import Path
from typing import Optional, List

from pptx import Presentation

from nannyml import ModelMetadata, FeatureType
from nannyml.performance_calculation import Metric


def generate_report(output_directory: Path, metadata: ModelMetadata, metrics: List[Metric] = None):
    date_string = datetime.now().strftime('%Y%m%d')
    prs = Presentation('report.pptx')

    # add title slide
    _add_title_slide(prs, model_name=metadata.name)

    _add_section_title(prs, section_name="Performance Estimation", section_subtext='CBPE / ROC AUC')

    performance_estimation_image_path = output_directory / 'performance_estimation/CBPE/images' / \
                                       'estimated_performance.png'
    if performance_estimation_image_path.exists():
        _add_image_slide(prs, performance_estimation_image_path, title="Estimated performance")

    _add_section_title(prs, section_name="Realized Performance")

    for metric in metrics:
        performance_calculation_image_path = output_directory / 'performance_calculation/images' / \
                                             f'realized_performance-{metric.display_name}.png'
        if performance_calculation_image_path.exists():
            _add_image_slide(prs, performance_calculation_image_path,
                             title=f"Realized performance: {metric.display_name}")

    # add drift stuff
    _add_section_title(prs, section_name="Model input drift", section_subtext='Univariate / Statistical')

    for feature in metadata.features:
        statistic_image_path = output_directory / 'drift/univariate/statistical/images' / \
                               f'statistical_drift-{feature.label}-statistic.png'
        if statistic_image_path.exists():
            statistic = 'Chi square statistic' if feature.feature_type == FeatureType.CATEGORICAL else 'KS statistic'
            _add_image_slide(prs, statistic_image_path, title=f"{statistic} for '{feature.label}'")

        p_values_image_path = output_directory / 'drift/univariate/statistical/images' / \
                              f'statistical_drift-{feature.label}-p_value.png'
        if p_values_image_path.exists():
            _add_image_slide(prs, p_values_image_path, title=f"P-values for '{feature.label}'")

        distribution_image_path = output_directory / 'drift/univariate/statistical/images' / \
                                  f'statistical_drift-{feature.label}-distribution.png'
        if distribution_image_path.exists():
            _add_image_slide(prs, distribution_image_path, title=f"Distribution for '{feature.label}'")

    _add_section_title(prs, section_name="Model input drift", section_subtext='Multivariate / Data Reconstruction')

    data_reconstruction_image_path = output_directory / 'drift/multivariate/data_reconstruction/images' / \
                                     'data_reconstruction_drift.png'
    if data_reconstruction_image_path.exists():
        _add_image_slide(prs, data_reconstruction_image_path, title="Data reconstruction")

    _add_section_title(prs, section_name="Model output drift", section_subtext='Univariate / Statistical')

    _add_section_title(prs, section_name="Target drift", section_subtext='Target distribution')

    prs.save(output_directory / f'{date_string}_nannyml_report.pptx')


def _add_title_slide(prs: Presentation, model_name: Optional[str]) -> Presentation:
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    slide.shapes.title.text = f'Report: {model_name}' if model_name else 'NannyML report'
    slide.shapes[1].text = f"Generated by NannyML - {datetime.now().strftime('%Y/%m/%d %H:%M')}"


def _add_section_title(prs, section_name: str, section_subtext: str = None) -> Presentation:
    slide = prs.slides.add_slide(prs.slide_layouts[8])

    slide.shapes.title.text = section_name
    slide.shapes[1].text = section_subtext if section_subtext else ''


def _add_image_slide(prs, image_path: Path, title: str = None) -> Presentation:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))

    slide.placeholders[13].insert_picture(str(image_path))
    slide.placeholders[0].text = title or ''
